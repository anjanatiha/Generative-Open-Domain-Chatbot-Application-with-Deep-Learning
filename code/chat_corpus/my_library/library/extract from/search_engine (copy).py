
# coding: utf-8

# In[171]:


'''
* Author           : Anjana Tiha
* Course           : Inform Retrieval/Web Search
* Course Instructor: Professor Visali Rus
* Semester         : Fall 2017
* University       : University of Memphis 
*
* Project Name     : Web Search Engine Implementation using Python(Implemented web crawler, preprocessor, inverted document
*                    indexer and cosine similarity measurer)  
* Last Update      : 12.17.2017
*
* Description      : 1. Implemented web search engine with query processing and web crawler with web document processor.
*                  : 2. Collected 10, 000 documents from "www.memphis.edu" web domain including webpages, .pdf, docx, pptx
*                       and .txt files.
*                  : 3. Implemented query processing engine by genereting TF-IDF for query and calculating cosine simmilarity
*                  :    between query and document corpus TF-IDF vector space.
*                  : 4. For document preprocessing of webpages/web documents including text files and pdf files:
*                       1. Removed HTML and script tags.
*                       2. Removing urls.
*                       3. Removing special characters and digits.
*                       4. Changing upper to lower case letters.
*                       5. Remove stop words.
*                       6. Stemming all words to root. 
*                    5. Built inverted document frequency index.
*                    6. Generated TF-IDF vector for full document corpus. 
*                    7. Calculates cosine similarity between query and  documents.
*                    8. Returns pages with maximum cosine similarity in decending order.
*                    9. Crawler can update and add new documents besides primary crawling.
*
* Tools Requirement: pdf2text (For .pdf processing)
*
* Comments         : Please use Anaconda editor for convenience.
* Version History  : 2.3.2.4
*                    Added incremental document collection option for web crawling an pre processing.
*                    New TF-IDF calculation is required.
* Current Version  : 2.3.4
*                    Changed search engine for working without re calculation for TF-IDF
'''


# In[172]:


### Import modules

import os, errno
import math
import time
import operator 
import collections
from collections import OrderedDict
from collections import deque
import queue
import shutil
import pickle
import re
import urllib
from urllib.parse import urlsplit
from urllib.parse import urljoin
import requests
from bs4 import BeautifulSoup, Tag
from nltk.stem.porter import PorterStemmer
import datetime
import numpy as np
import docx2txt
from pptx import Presentation
import copy


# In[173]:



# global Hashmaps, variables
stopwords = {}
term_doc_freq_vector = {}
doc_term_freq_vector = {}
doc_term_freq_vector_norm = {}
page_doc_map = {}
doc_page_map = {}
page_ref_count = {}
doc_count = 0
link_queue = queue.Queue()
last_doc_index = -1
page_queued_map = {}
start_time = time.time()
total_number_docs = 0



# In[174]:


#Global directory or file name

crawled_web_dir = "web_text_crawled"
crawled_web_dir_conv_need = "web_docs_crawled"
crawled_web_dir_preprocessed = "web_text_preprocessed"
output_web_dir = "output"

stopword_path = "english.stopwords.txt"

list_dir = [crawled_web_dir, crawled_web_dir_conv_need, crawled_web_dir_preprocessed]

url = "http://www.memphis.edu/"
#url = "http://www.cs.memphis.edu/~vrus/teaching/ir-websearch/"
domain = "memphis.edu"
total_number_docs = 10000


# In[175]:


# File Operations
# Create/Delete/ file/Directory 

# create one single directory
def create_directory(directory):
    try:
        os.makedirs(directory)
    except OSError as e:
        if(e.errno != errno.EEXIST):
            raise
    pass 


# create a list of directories
def create_directories(list_dir):
    for dir_i in list_dir:
        print(dir_i)
        create_directory(dir_i) 
        
        
# delete one single directory
def delete_directory(dir_name):
    if(os.path.isdir(dir_name)):
        try:
            shutil.rmtree(dir_name)
        except OSError as e:
            if(e.errno != errno.EEXIST):
                raise
        pass 

    
# delete a list of directories
def delete_directories(list_dir):
    for dir_i in list_dir:
        print(dir_i)
        delete_directory(dir_i)

        
#delete if file is empty
def delete_file(path):
    try:
        os.remove(path)
    except WindowsError:
        print("failed deleting: " + path)
        pass


# In[176]:



#save text in given path with given file name
def save_text(text, dir_path, file_name):
    text_file = open(dir_path+"\\"+file_name, "w")
    text_file.write(text)
    text_file.close()
        
        
#load stop words from given file path
def load_stopwords(filepath):
    with open(filepath, 'r') as content_file:
        for line in content_file:
            line = line.strip()
            stopwords[line] = 1
            


# In[177]:



# save object in pickle
def save_obj(obj, name, key_or_val, order):
    filename = name + ".p"
    
    if(key_or_val == "key" and order == "auto"):
        sorted_x = sorted(obj.items(), key=operator.itemgetter(0))
    elif(key_or_val == "key" and order == "reverse"):
        sorted_x = sorted(obj.items(), key=operator.itemgetter(0), reverse=True)
    elif(key_or_val == "value" and order == "auto"):
        sorted_x = sorted(obj.items(), key=operator.itemgetter(1))
    elif(key_or_val == "value" and order == "reverse"):
        sorted_x = sorted(obj.items(), key=operator.itemgetter(1), reverse=True)
    if (os.path.isfile(filename)):
        os.remove(filename)
        
    pickle.dump( obj, open( filename, "wb" ) )


#save object     
def save_obj_without_sort(obj, name):
    pickle.dump( obj, open( name + ".p", "wb" ) )

    
# save object in pickle
def save_obj_no_sort(obj, name):
    filename = name + ".p"
    
    if (os.path.isfile(filename)):
        os.remove(filename)
        
    pickle.dump( obj, open( filename, "wb" ) )
    
    
#save current queue for future resume of crawling
def save_obj_no_sort_w(queue1, filename):
    link_list = []
    
    new_queue = queue.Queue()
    new_queue.queue = copy.deepcopy(queue1.queue)
    if (os.path.isfile(filename)):
        os.remove(filename)
        
    while(new_queue.empty() == False):
        link_list.append(new_queue.get())
        
    filename = filename + ".p"
    
    pickle.dump( link_list, open( filename, "wb" ) )


    
#load object from pickle file
def load_obj(name):
    file = open(name,'rb')
    object_file = pickle.load(file)
    file.close()
    
    return object_file
  
    
# save object in pickle
def load_obj_no_sort(name):
    file = open(name,'rb')
    object_file = pickle.load(file)
    file.close()
    
    return object_file
 
    
#load last visit queue for resume of crawling
def load_obj_no_sort_w(filename):
    global link_queue
    link_list =[]
    
    if (os.path.isfile(filename)):
        file = open(filename,'rb')
        link_list = pickle.load(file)
        file.close()
        
        for link in link_list:
            link_queue.put(link)
        return link_queue
    
    else:
        print("no file found")
        return

#reset all variables
def reset_global_variables():
    global stopwords
    global term_doc_freq_vector
    global doc_term_freq_vector
    global page_doc_map
    global doc_page_map
    global page_ref_count
    global doc_count
    global link_queue
    global last_doc_index
    global page_queued_map
    global start_time

    stopwords = {}
    term_doc_freq_vector = {}
    doc_term_freq_vector = {}
    page_doc_map = {}
    doc_page_map = {}
    page_ref_count = {}
    doc_count = 0
    link_queue = queue.Queue()
    last_doc_index = -1
    page_queued_map = {}
    start_time = time.time()
    
    
#delete all files
def delete_all_files():
    try:
        delete_directories(list_dir)
    except:
        pass
    try:
        delete_file("doc_count.p")
    except:
        pass
    try:
        delete_file("doc_term_freq_vector.p")
    except:
        pass
    try:
        delete_file("doc_term_freq_vector_norm.p")
    except:
        pass
    try:
        delete_file("doc_url_map.p")
    except:
        pass
    try:
        delete_file("link_queue.p")
    except:
        pass
    try:
        delete_file("page_ref_count.p")
    except:
        pass
    try:
        delete_file("term_doc_freq_vector.p")
    except:
        pass
    try:
        delete_file("url_doc_map.p")
    except:
        pass

    
# save crawled doc object    
def save_all_obj():
    global page_doc_map
    global doc_page_map
    global page_ref_count
    global doc_count
    global link_queue
    
    save_obj(page_doc_map, "url_doc_map", "value", "auto")
    save_obj(doc_page_map, "doc_url_map", "key", "auto")
    save_obj_no_sort(doc_count, "doc_count")
    save_obj_no_sort_w(link_queue, "link_queue")
    save_obj(page_ref_count, "page_ref_count", "value", "reverse")
    print("saved objects")

#save tfidf    
def save_all_obj_tfidf(doc_term_freq_vector_norm_new):
    global term_doc_freq_vector
    global doc_term_freq_vector
    
    save_obj_without_sort(term_doc_freq_vector, "term_doc_freq_vector")
    save_obj_without_sort(doc_term_freq_vector, "doc_term_freq_vector")
    save_obj_without_sort(doc_term_freq_vector_norm_new, "doc_term_freq_vector_norm")
    
    
    
# load crawled doc object  
def load_all_obj():
    global page_doc_map
    global doc_page_map
    global page_ref_count
    global doc_count
    global link_queue
    
    page_doc_map = load_obj("url_doc_map.p")
    doc_page_map = load_obj("doc_url_map.p")
    #doc_count = load_obj_no_sort("doc_count.p")
    doc_count = max(doc_page_map.keys())
    link_queue = load_obj_no_sort_w("link_queue.p")
    page_ref_count = load_obj("page_ref_count.p")
    print("loaded objects")


def load_obj_search():
    global total_number_docs
    global doc_url_map
    global term_doc_freq_vector
    global doc_term_freq_vector
    global doc_term_freq_vector_norm
    
    stopword_path = "english.stopwords.txt"
    doc_url_map_file = "doc_url_map.p"
    term_doc_freq_file = "term_doc_freq_vector.p"
    doc_term_freq_file = "doc_term_freq_vector.p"
    doc_term_freq_file_norm = "doc_term_freq_vector_norm.p"
    doc_count = "doc_count.p"
    
    total_number_docs = load_obj(doc_count)
    doc_url_map = load_obj(doc_url_map_file)
    term_doc_freq_vector = load_obj(term_doc_freq_file)
    doc_term_freq_vector = load_obj(doc_term_freq_file)
    doc_term_freq_vector_norm = load_obj(doc_term_freq_file_norm)
    load_stopwords(stopword_path)

    


# In[178]:


#print elapsed time in hh:mm:ss format
def format_time(start_time, end_time):
    elsapsed_time = end_time - start_time
    hr = int(elsapsed_time)//3600
    min_ = (int(elsapsed_time) - (hr * 3600))/60
    sec = int(elsapsed_time) - hr * 3600 - min_ * 60
    print("HH:Min:Sec > " + str(hr) +" hr " + str(min_) + " min "+ str(sec) + "sec")



# In[179]:



# URLS and Text preprocessing functions

# remove fragment identifier # and repeated loop url for php and asp
def remove_url_frag_id(url):
    if ".php" in url:
        url = url.split('.php')
        if(len(url)>1):
            url =url[0] + ".php"
    elif ".aspx" in url:
        url = url.split('.aspx')
        if(len(url)>1):
            url =url[0] + ".aspx"   
    url = url.split('#')[0]
    
    return url


# remove fragment identifier #
def remove_url_frag_simple(url):   
    url = url.split('#')[0]
    return url


# removes "/" from url
def remove_slash_before_or_after(url, type_r):
    if(type_r == "before"):
        if url.startswith("/"):
            url = url[1:]
        return url
    
    elif(type_r == "after"):   
        if url[-1]=="/":
            url = url.rsplit('/', 1)[0]
        return url 
    
    
# remove http or https from webpage urls
def strip_http_s(url):    
    url = url.replace("https://","")
    url = url.replace("http://","")
    url = url.rstrip('\/') 
    url = "http://"+ url
    
    return url


# check if url is in selected domain name
def check_if_in_domain(url, domain):
    if(domain in url): 
        return 1
    else:
        return 0

    
# file type extention type exclusion
def is_excluded_type(extension):
    exclude_list = ["jpg", "jpeg", "png", "mp3", "mp4", "xlx"]
    if extension in exclude_list:
        return 1
    else:
        return 0

    
#check whether URL is valid
def check_valid_URL(url):
    url_reg = re.compile(
        r'^(?:http|ftp)s?://' # http:// or https://
        r'(?:(?:[A-Z0-9](?:[A-Z0-9-]{0,61}[A-Z0-9])?\.)+(?:[A-Z]{2,6}\.?|[A-Z0-9-]{2,}\.?)|' #domain...
        r'localhost|' #localhost...
        r'\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3})' # ...or ip
        r'(?::\d+)?' # optional port
        r'(?:/?|[/?]\S+)$', re.IGNORECASE)
    
    is_valid = url_reg.match(url)
    
    return is_valid


#get extention of link to check the link type(.txt, .pdf, or html) 
def get_page_extention(url):
    weblink_extention = url.rsplit('.', 1)[-1]
    return weblink_extention


#remove hyperlink from web page text for preprocessing
def remove_hyper_link(text):
    URLless_string = re.sub(r'(?i)\b((?:https?://|www\d{0,3}[.]|[a-z0-9.\-]+[.][a-z]{2,4}     /)(?:[^\s()<>]+|\(([^\s()<>]+|(\([^\s()<>]+\)))*\))+(?:\(([^\s()<>]+|(\([^\s()<>]+\)))*\)|[^\s`!()\[\]{};:\'".,<>?«»“”‘’]))', '', text)
    return URLless_string


#remove special character in single line
def remove_special_char(line):
    line = re.sub('[^a-zA-Z]+', ' ', line)
    return line


# get all the url/links to other pages from current pages        
def get_all_links(url, html):
    global domain
    global link_queue
    global page_queued_map
    
    soup = BeautifulSoup(html, "html.parser")
    links = soup.find_all('a')
    
    
    for tag in links:
        link = tag.get('href', None)
        
        if link is not None:
            try:
                link_extention = get_page_extention(link)
                
                if(link == "" or link == "#" or link_extention == "ppt" or is_excluded_type(link_extention) == 1):
                    a=1

                elif(link_extention in ["pdf", "docx", "pptx", "txt"]):
                    if(check_valid_URL(link)):
                        link_original = strip_http_s(link)
                        
                        if link_original not in page_queued_map:
                            page_queued_map[link_original] = 1
                            link_queue.put(link)

                    else:
                        modified_url = remove_url_frag_id(url)
                        modified_url = remove_slash_before_or_after(modified_url, "after")
                        modified_link = remove_slash_before_or_after(link, "before")
                        modified_link = modified_url + "/" + modified_link
                        
                        if(check_valid_URL(modified_link)):
                            link_original = strip_http_s(modified_link)
                            
                            if link_original not in page_queued_map:
                                page_queued_map[link_original] = 1
                                link_queue.put(modified_link)

                else:
                    is_valid = check_valid_URL(link)
                    if(is_valid):
                        modified_link = remove_url_frag_id(link)
                        modified_link = remove_slash_before_or_after(modified_link, "after")
                        link_original = strip_http_s(modified_link)
                        
                        if(check_if_in_domain(modified_link, domain) == 1):
                            
                            if link_original not in page_queued_map:
                                page_queued_map[link_original] = 1
                                link_queue.put(modified_link)

                    else:
                        modified_url= remove_url_frag_id(url)
                        modified_url = remove_slash_before_or_after(modified_url, "after")
                        modified_link = remove_slash_before_or_after(url, "before")
                        
                        if(modified_url!=modified_link):
                            
                            if  modified_url not in modified_link:
                                modified_link = modified_url + "/" + modified_link 
                                
                                if(check_if_in_domain(modified_link, domain) == 1):
                                    if(check_valid_URL(modified_link)):
                                        link_original = strip_http_s(modified_link)
                                        
                                        if link_original not in page_queued_map:
                                            page_queued_map[link_original] = 1
                                            link_queue.put(modified_link) 
                            else: 
                                if(check_if_in_domain(modified_link, domain) == 1):
                                    if(check_valid_URL(modified_link)):
                                        link_original = strip_http_s(modified_link)
                                        
                                        if link_original not in page_queued_map:
                                            page_queued_map[link_original] = 1
                                            link_queue.put(modified_link)
            except:
                continue


# In[180]:



# convert pdf to text using "pdf2text"
def pdf_to_text(input_pdf, file_name):
    global crawled_web_dir
    os.system(("pdftotext %s %s") %( input_pdf, crawled_web_dir+"//"+file_name))

    
# convert pptx to text    
def pptx_to_text(book_path):
    prs = Presentation(book_path)
    text = ""
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = text_runs + run.text
    return text


#import pdf, docx, pptx from url/single web link and convert to text and save in directory
def import_convert_preprocess(url, extension):
    global doc_count
    global crawled_web_dir_preprocessed
    global crawled_web_dir
    global crawled_web_dir_conv_need
    global page_doc_map
    url_map_name = url
    
    if(url_map_name not in page_doc_map):
        page_doc_map[url_map_name] = -1
        page_ref_count[url_map_name] = 1
        
        try:
            doc_count_temp = doc_count + 1
            book_name = ""
            if extension == "pdf":
                book_name = str(doc_count_temp) + ".pdf"
            elif extension == "docx":
                book_name = str(doc_count_temp) + ".docx"
            elif extension == "pptx":
                book_name = str(doc_count_temp) + ".pptx"

            book_path = crawled_web_dir_conv_need + "\\" + book_name
            
            a = requests.get(url, stream=True)
            
            with open(book_path, 'wb') as book:   
                for block in a.iter_content(512):
                    if not block:
                        break
                    book.write(block)
                    
            book.close()

            file_name = str(doc_count_temp)+".txt"
            file_path = crawled_web_dir+ "\\" + file_name
            is_valid_for_indexing = 555
            if extension == "pdf":
                pdf_to_text(book_path, file_name)
                is_valid_for_indexing = preprocess_one_doc_from_pdf(crawled_web_dir, file_name, crawled_web_dir_preprocessed)
                
            elif extension == "docx":
                text = docx2txt.process(book_path)
                save_text(text, crawled_web_dir, file_name)
                is_valid_for_indexing = preprocess_one_doc(crawled_web_dir, file_name, crawled_web_dir_preprocessed)
                
            elif extension == "pptx":
                text = pptx_to_text(book_path)
                save_text(text, crawled_web_dir, file_name)
                is_valid_for_indexing = preprocess_one_doc(crawled_web_dir, file_name, crawled_web_dir_preprocessed)
        

            if(is_valid_for_indexing == 1) :
                doc_count = doc_count + 1
                page_doc_map[url_map_name] = doc_count
                doc_page_map[doc_count] = url_map_name
                page_ref_count[url_map_name] = 1
            else:
                delete_file(book_path)
                delete_file(file_path)
                page_doc_map[url_map_name] = -2

        except IOError:
            page_doc_map[url_map_name]= -1
    else:
        page_ref_count[url_map_name] = page_ref_count[url_map_name] + 1

        

#remove all html and scripting
def remove_extra_space(txt):
    # Removes all blank lines
    txt = re.sub(r'\n\s*\n', '\n', txt)
    return txt


# clean html of the tag and markups
def clean_html(html_text):
    global crawled_web_dir
    soup = BeautifulSoup(html_text, "html.parser")

    for script in soup(['style', 'script', 'head', 'title', 'meta', '[document]']):
        script.extract()
    for tag in soup.find_all('a'):
        tag.replaceWith('')
    for tag in soup.find_all('footer'):
        tag.replaceWith('')
    
    clean_text = soup.get_text()
    clean_text = remove_extra_space(clean_text)
    return clean_text




# import text from single web page and preprocess
def fetch_extract_html_txt(url):
    global doc_count
    global domain
    global crawled_web_dir
    global crawled_web_dir_preprocessed
    
    if(check_if_in_domain(url, domain) == 0):
        return 
    url_map_name = url
    
    if(url_map_name in page_doc_map):
        page_ref_count[url_map_name] = page_ref_count[url_map_name] + 1
        
    else:
        page_doc_map[url_map_name] = -1
        page_ref_count[url_map_name] = 1
        
        try: 
            html = urllib.request.urlopen(url) 
            html_text = html.read() 
            
            if(html_text.strip() == ""):
                return
            
            clean_text = clean_html(html_text)
            clean_text = clean_text.strip()

            if clean_text.strip()=="":
                return
            
            doc_count = doc_count + 1
            page_doc_map[url_map_name] = doc_count
            doc_page_map[doc_count] = url_map_name
            save_text(clean_text, crawled_web_dir, str(doc_count)+".txt")
            
            file_name = str(doc_count)+".txt"
            file_path = crawled_web_dir+ "\\" + file_name
            
            is_valid_for_indexing = preprocess_one_doc(crawled_web_dir, file_name, crawled_web_dir_preprocessed)
            
            if(is_valid_for_indexing == 0) :
                delete_file(file_path)
                page_doc_map[url_map_name] = -2
                doc_count = doc_count - 1
                
            get_all_links(url, html_text)
            
        except:
            page_doc_map[url_map_name]= -1
            



# In[181]:


# preprocess files in a folder to remove punctuations, digits, special characters, url/web links
# removes stop words given in file
# convert to origin word/ do stemming


# preprocess webpage text, docx, pptx text and save in directory
def preprocess_one_doc(input_dir, input_filename, output_dir):
    ps = PorterStemmer()
    input_file_path = input_dir + "\\"+ input_filename
    text = ""
    count = 0
    
    try:
        with open(input_file_path, 'r') as content_file:
            for line in content_file:
                if(line in ['\n', '\r\n','\r']):
                    continue
                    
                line = line.strip()
                line = remove_hyper_link(line)
                line = remove_special_char(line)
                line = line.lower()
                line = re.sub(' +',' ',line)
                words = line.split(" ")
                
                for word in words:
                    word = word.strip()
                    word = remove_special_char(word)
                    word = re.sub(' +','',word)
                    
                    if word not in stopwords and word != " " and word != "":
                        stem_word = ps.stem(word)
                        text = text + " " + stem_word
                        count = count + 1       
                        
            if(count > 50):
                save_text(text, output_dir, input_filename)
                return 1
            
            else:      

                return 0
    except:
        return 0
        
    
    
    
# preprocess pdf text and  save in directory
def preprocess_one_doc_from_pdf(input_dir, input_filename, output_dir):
    ps = PorterStemmer()
    input_file_path = input_dir + "\\"+ input_filename
    text = ""
    count = 0
    ret_val = 999
    
    try:
        with open(input_file_path, 'rb') as content_file:
            for line in content_file:
                line = line.decode("utf-8")
                if(line in ['\n', '\r\n','\r']):
                    continue
                    
                line = line.strip()
                line = remove_hyper_link(line)
                line = remove_special_char(line)
                line = line.lower()
                line = re.sub(' +',' ',line)
                words = line.split(" ")

                for word in words:
                    word = word.strip()
                    word = remove_special_char(word)
                    word = re.sub(' +','',word)
                    
                    if word not in stopwords and word != " " and word != "":
                        stem_word = ps.stem(word)
                        text = text + " " + stem_word
                        count = count + 1
                        
            if(count > 50):
                save_text(text, output_dir, input_filename)
                return 1
            
            else:
                return  0
    except:
        return 0



# In[182]:



# crawl through single webpage
def webpage_crawler(total_number_docs):
    global doc_count 
    global link_queue
    global last_doc_index
    global start_time
    global crawled_web_dir
    global crawled_web_dir_conv_need
    
    if(doc_count % 100 == 0 and last_doc_index != doc_count):
        print("Extracted Documents: " + str(doc_count))
        last_doc_index = doc_count
        
    if(doc_count % 200 == 0):
        format_time(start_time, time.time())
        
    url = link_queue.get()
    #print(doc_count+1, " : ", url)
    
    try:
        link_extention = get_page_extention(url)
        
        if(url == "" or link_extention == "ppt"):
            a=1
        elif(link_extention in ["pdf","docx", "pptx"]):
            import_convert_preprocess(url, link_extention) 
        elif(link_extention == "txt"):
            fetch_extract_html_txt(url)
        else:
            fetch_extract_html_txt(url)
            
    except: 
        pass        


    
# crawl through a website            
def website_crawler(total_number_docs):
    global link_queue
    global doc_count
    global total_number_doc
    
    while(doc_count < total_number_docs):
        if(link_queue.empty()):
            print("Queue is empty")
            return
        
        if(doc_count%200 == 0 and doc_count != 0):
            save_all_obj()
            
        webpage_crawler(total_number_docs)
        
    save_all_obj()


# In[183]:



#web crawler main
def web_crawling_main(url, domain, total_page_count):
    print("Start Time: ", datetime.datetime.time(datetime.datetime.now()))
    
    delete_all_files()
    reset_global_variables()
    
    create_directories(list_dir)
    
    load_stopwords(stopword_path)
    
    url = remove_slash_before_or_after(url, "after")
    link_original = strip_http_s(url)
    page_queued_map[link_original] = 1
    print(url)
    link_queue.put(url)
    
    total_number_docs = total_page_count
    
    website_crawler(total_number_docs)
    

# Update crawling from last saved url location    
def web_crawling_main_update(url, domain, num_add_doc):
    global page_doc_map
    global doc_page_map
    global page_ref_count
    global total_number_docs
    global doc_count
    
    print("Start Time: ", datetime.datetime.time(datetime.datetime.now()))
    
    reset_global_variables()
    load_stopwords(stopword_path)
    load_all_obj()

    total_number_docs = doc_count + num_add_doc
    website_crawler(total_number_docs)

    


# In[184]:


def cosine_similarity(v1,v2):
    sumxx, sumxy, sumyy = 0, 0, 0
    for i in range(len(v1)):
        x = v1[i]; 
        y = v2[i]
        sumxx += x*x
        sumyy += y*y
        sumxy += x*y
    return sumxy/math.sqrt(sumxx*sumyy)


# In[185]:



#build inverted index for all files present in preprocessed file directory
def inverse_document_indexer(preprocessed_file_dir_path):
    dirs = os.listdir(preprocessed_file_dir_path)
    i = 0 
    
    for file in dirs:
        filepath = preprocessed_file_dir_path + "\\"+ file
        text = ""
        i = i + 1
        
        if(i % 1000 == 0):
            print("Building inverse document index for file no: "+str(i))   
            print("Current Time: ", datetime.datetime.time(datetime.datetime.now()))
            
        try:
            with open(filepath, 'r') as content_file:
                file_name = str(file)[:-4]
                
                doc_term_freq_vector[file_name] = {}
                single_doc_term_freq_vector = doc_term_freq_vector[file_name]
                
                for line in content_file:
                    line = line.strip()
                    words = line.split(" ")
                    
                    for word in words:
                        word = word.strip()
                        
                        if word != "":
                            if word not in term_doc_freq_vector:
                                single_term_doc_freq_vector = {}
                                single_term_doc_freq_vector[file_name] = 1
                                single_term_doc_freq_vector["DocFreq"] = 1  
                                term_doc_freq_vector[word] = single_term_doc_freq_vector

                            else:

                                single_term_doc_freq_vector = term_doc_freq_vector[word]

                                if file_name not in single_term_doc_freq_vector:
                                    single_term_doc_freq_vector[file_name] = 1
                                    single_term_doc_freq_vector["DocFreq"] = single_term_doc_freq_vector["DocFreq"] + 1
                                    term_doc_freq_vector[word] = single_term_doc_freq_vector

                                else:
                                    single_term_doc_freq_vector[file_name] = single_term_doc_freq_vector[file_name] + 1
                                    term_doc_freq_vector[word] = single_term_doc_freq_vector 

                            a=1
                            if "DocMaxFreq" not in single_doc_term_freq_vector:
                                single_doc_term_freq_vector["DocMaxFreq"] = 1

                            if word not in single_doc_term_freq_vector:
                                single_doc_term_freq_vector[word] = 1
                                doc_term_freq_vector[file_name] = single_doc_term_freq_vector
                                
                            else:
                                single_doc_term_freq_vector[word] = single_doc_term_freq_vector[word] + 1
                                
                                if(single_doc_term_freq_vector[word] > single_doc_term_freq_vector["DocMaxFreq"]):
                                    single_doc_term_freq_vector["DocMaxFreq"] = single_doc_term_freq_vector[word]
                                doc_term_freq_vector[file_name] = single_doc_term_freq_vector

        except:
            pass

           
            
#builds tfidf from inverse document index
def tfidf_document_text(term_doc_freq_vector, doc_term_freq_vector):
    global total_number_docs
    
    total_number_docs = load_obj("doc_count.p")
    
    doc_term_freq_vector_normalized = doc_term_freq_vector

    
    for doc in doc_term_freq_vector_normalized:
        for term in doc_term_freq_vector_normalized[doc]:
            
            if(term != "DocMaxFreq"):
                doc_freq = term_doc_freq_vector[term]["DocFreq"]
                doc_term_freq_vector_normalized[doc][term] = (doc_term_freq_vector[doc][term]/doc_term_freq_vector[doc]["DocMaxFreq"])*(math.log2(total_number_docs/doc_freq))
    
    
    for doc in doc_term_freq_vector_normalized:
        del doc_term_freq_vector_normalized[doc]["DocMaxFreq"]
        
    return doc_term_freq_vector_normalized



def inverse_document_indexer_final(crawled_web_dir_preprocessed, stopword_path):
    inverse_document_indexer(crawled_web_dir_preprocessed)
    doc_term_freq_vector_norm = tfidf_document_text(term_doc_freq_vector, doc_term_freq_vector)
    save_all_obj_tfidf(doc_term_freq_vector_norm)
    
    return term_doc_freq_vector, doc_term_freq_vector, doc_term_freq_vector_norm


# In[186]:


#query preprocessing 
def query_preprocessor(query_str):
    ps = PorterStemmer()
    query_dict = {}
    
    query_str_modified = query_str.strip()
    query_str_modified = remove_special_char(query_str_modified)
    query_str_modified = query_str_modified.lower()
    query_str_modified = re.sub(' +',' ',query_str_modified)
    words = query_str_modified.split(" ")
    
    max_freq = 0
    N = 0
    
    for word in words:
        word = word.strip()
        
        if word not in stopwords and word !="": 
            word = ps.stem(word)
            
            if word not in query_dict:
                query_dict[word] = 1
            else:
                query_dict[word] = query_dict[word] + 1
                
            if(query_dict[word] > max_freq):
                max_freq = query_dict[word]
                
            N +=1
            
    return query_dict, max_freq, N


#generate normalized term vector for query
def query_normalizer(query_dict, max_freq, total_number_docs, term_doc_freq_vector):
    query_dict_normalized = {}
    doc_term_freq_vector_normalized = doc_term_freq_vector 
    
    for word in query_dict:
        if word in term_doc_freq_vector:
            query_dict_normalized[word] =  ( 0.5  +  (0.5 * query_dict[word] / max_freq) ) * (math.log2((total_number_docs+1)/(term_doc_freq_vector[word]["DocFreq"]+1)))
        else:
            query_dict_normalized[word] =  ( 0.5  +  (0.5 * query_dict[word] / max_freq) ) * (math.log2((total_number_docs+1)))
    
    return query_dict_normalized




# retrieve documents for query
def retrieve_docs_with_query_word(query_term_freq_vect_norm, term_doc_freq_vector):
    docs_with_query_terms = []
    
    for word in query_term_freq_vect_norm:
        if word in term_doc_freq_vector:
            docs = term_doc_freq_vector[word]
            for doc in docs:
                if(doc != "DocFreq"):
                    docs_with_query_terms.append(doc)
                    
    return docs_with_query_terms

# retrieve relevant document
def calculate_cosine_query_doc(docs_with_query_terms, query_term_freq_norm, term_doc_freq_vector, doc_term_freq_vector):
    cosine_query_doc = {}
    
    for doc in docs_with_query_terms:
        temp = {}
        
        for word in doc_term_freq_vector_norm[doc]:
            if word in query_term_freq_norm:
                temp[word] = query_term_freq_norm[word]
            else:
                temp[word] = 0
                
        doc_v=[]
        query_v = []
        
        for word in doc_term_freq_vector_norm[doc]:
            doc_v.append(doc_term_freq_vector_norm[doc][word])
            query_v.append(temp[word])
        cosine_query_doc[doc] = cosine_similarity(doc_v, query_v)
    
    return cosine_query_doc


#get url from file namedoc
def get_url(cosine_query_doc, doc_url_map):
    url_list = []
    similarity = []
    doc_list = []
    similarity_map = {}
    
    cosine_query_doc_new = sorted(cosine_query_doc.items(), key=operator.itemgetter(1), reverse = True)  
    cosine_query_doc_newest = {}
    
    for doc in cosine_query_doc_new:
        cosine_query_doc_newest[doc[0]] = doc[1]
        
    for doc in cosine_query_doc_newest:
        url_list.append(doc_url_map[int(doc)])
        similarity.append(cosine_query_doc_newest[doc])
        similarity_map[doc] = cosine_query_doc_newest[doc]
        doc_list.append(doc)
        
    return url_list, doc_url_map, similarity, similarity_map, doc_list
    


# In[187]:


#Evaluation
#finds relevant docs
def relevant_doc(query_term_freq_vect, doc_term_freq_vector):
    relevant_list = []
    relevant_list_map = {}
    
    for doc in doc_term_freq_vector:
        doc_i = 0
        
        for term in query_term_freq_vect:
            if term not in doc_term_freq_vector[doc]:
                doc_i = -1
                break
                
            else:
                doc_i = doc_i + 1
                
        if(doc_i!=-1):
            relevant_list.append(doc)
            relevant_list_map[doc] = doc_i
            
            
    return len(relevant_list), relevant_list, relevant_list_map


#calculate number of relevant docs
def num_relevant_doc_in_query(doc_list, query_term_freq_vect, doc_term_freq_vector):
    relevant_list = []
    relevant_list_map = {}
    
    for doc in doc_list:
        doc_i = 0
        
        for term in query_term_freq_vect:
            if term not in doc_term_freq_vector[str(doc)]:
                doc_i = -1
                break
                
            else:
                doc_i = doc_i + 1
                
        if(doc_i!=-1):
            relevant_list.append(doc)
            relevant_list_map[doc] = doc_i
            
    return len(relevant_list), relevant_list, relevant_list_map


# find precision, recall and f1 score
def evaluation(num, relevant_list_len, qrelevant_list_len):
    print(num, relevant_list_len, qrelevant_list_len)
    recall = -999
    precision = -999
    f1 = -999
    
    if relevant_list_len != 0:
        recall = qrelevant_list_len/relevant_list_len
    
    if num != 0:
        precision = qrelevant_list_len/num
        
    if recall != -999 or precision != -999:
        if recall == -999:
            recall = 0
        elif precision == -999 :
            precision = 0
            
        f1= (2*precision*recall)/(precision + recall)
    else:
        precision = 0
        recall = 0
        f1 = 0
    
    return precision, recall, f1


def perfomance(query_str, num):
    avg_precision = 0
    avg_recall = 0
    avg_f1 = 0
    query_str_len = len(query_str)
    i=1
    
    for query in query_str:
        url_list, doc_url_map, similarity, similarity_map, docs_with_query_terms, term_doc_freq_vector, query_term_freq_vect,doc_term_freq_vector,doc_list = web_search_main(query_str[query])
        relevant_list_len, relevant_list, relevant_list_map = relevant_doc(query_term_freq_vect, doc_term_freq_vector)
        qrelevant_list_len,relevant_list, qrelevant_list_map = num_relevant_doc_in_query(doc_list[:num], query_term_freq_vect, doc_term_freq_vector)
        
        precision, recall, f1 = evaluation(num, relevant_list_len, qrelevant_list_len)
        
        avg_precision = avg_precision + precision
        avg_recall = avg_recall + recall
        avg_f1 = avg_f1 + f1
        
        print(i," : ", query_str[query], " : precision : ", precision, ", recall : ", recall, ", f1 : ", f1)
        
        i = i + 1
        
        
    avg_precision = (avg_precision/query_str_len)
    avg_recall = (avg_recall/query_str_len)
    avg_f1 = (avg_f1/query_str_len)
    print("Average precision : ", avg_precision)
    print("Average recall : ", avg_recall)
    print("Average f1 : ", avg_f1) 
    
    
    return avg_precision, avg_recall, avg_f1


# In[195]:


# main search function 
def web_search_main(query_str):
    #time_1 = time.time()
    global total_number_docs
    global doc_url_map
    global term_doc_freq_vector
    global doc_term_freq_vector
    global doc_term_freq_vector_norm
    
    load_obj_search()
    
    query_term_freq_vect, max_freq, N = query_preprocessor(query_str)
    query_term_freq_vect_norm = query_normalizer(query_term_freq_vect, max_freq, total_number_docs, term_doc_freq_vector)
    docs_with_query_terms = retrieve_docs_with_query_word(query_term_freq_vect_norm, term_doc_freq_vector)
    cosine_query_doc = calculate_cosine_query_doc(docs_with_query_terms, query_term_freq_vect_norm, term_doc_freq_vector, doc_term_freq_vector)
    url_list, doc_url_map, similarity, similarity_map, doc_list = get_url(cosine_query_doc, doc_url_map)
    #format_time(time_1, time.time())
    
    return url_list, doc_url_map, similarity, similarity_map, docs_with_query_terms, term_doc_freq_vector, query_term_freq_vect, doc_term_freq_vector,doc_list


def search_engine_final_main(query_str, count):
    i = 0
    result = []
    url_list, doc_url_map, similarity, similarity_map, docs_with_query_terms, term_doc_freq_vector, query_term_freq_vect, doc_term_freq_vector,doc_list = web_search_main(query_str) 
    
    for url in url_list:
        url_row = []
        print(i+1, ". ", url, "\nSimillarity: ", similarity[i])
        url_row.append(str(i + 1)+".")
        url_row.append(url)
        url_row.append(similarity[i])
        result.append(url_row)
        i += 1
        
        if i>count:
            break
            
    return result


# In[196]:


#total_page_count = 9000

#web_crawling_main(url, domain, total_page_count)


# In[197]:



#num_add_doc = 200
#web_crawling_main_update(url, domain, num_add_doc)


# In[198]:


crawled_web_dir_preprocessed = "web_text_preprocessed"
stopword_path = "english.stopwords.txt"
#term_doc_freq_vector, doc_term_freq_vector, doc_term_freq_vector_norm = inverse_document_indexer_final(crawled_web_dir_preprocessed, stopword_path)


# In[199]:


query_str = {"q1":"international office", "q2":"software engineering research", "q3":"Cookie", "q4":"president of the university","q5":"computer science research awards",
"q6":"semantic similarity", "q7":"tiger bike's current offer","q8": "where is the library located?", "q9":"How to graduate with honors?","q10":"scholarships in computer science"}


# In[200]:


#num=100
#perfomance(query_str, num)

